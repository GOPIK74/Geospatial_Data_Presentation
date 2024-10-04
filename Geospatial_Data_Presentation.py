from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
ppt = Presentation()

# Title Slide
slide_1 = ppt.slides.add_slide(ppt.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]

title.text = "Exploring Digital Elevation Models (DEM) and Satellite Data for Geospatial Projects"
subtitle.text = "2024 Available DEM Data and Satellites"

# Slide 2 - SRTM
slide_2 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_2.shapes.title
content = slide_2.placeholders[1]

title.text = "#SRTM (Shuttle Radar Topography Mission)"
content.text = (
    "Resolution: 30m, 90m\n"
    "Orbit Rotation: 90 minutes (Low Earth Orbit)\n"
    "Distance from Earth: ~233 km\n"
    "Availability: Global, via USGS Earth Explorer.\n"
    "Year: Launched in 2000.\n"
    "Price: Free."
)

# Slide 3 - ASTER GDEM
slide_3 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_3.shapes.title
content = slide_3.placeholders[1]

title.text = "#ASTER GDEM (Advanced Spaceborne Thermal Emission and Reflection Radiometer)"
content.text = (
    "Resolution: 30m\n"
    "Orbit Rotation: 99 minutes\n"
    "Distance from Earth: 705 km\n"
    "Availability: Global, via NASA Earthdata.\n"
    "Year: Launched in 1999.\n"
    "Price: Free."
)

# Slide 4 - Copernicus DEM
slide_4 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_4.shapes.title
content = slide_4.placeholders[1]

title.text = "#Copernicus DEM (Sentinel Mission)"
content.text = (
    "Resolution: 30m, 90m\n"
    "Orbit Rotation: ~98.6 minutes (Polar orbit)\n"
    "Distance from Earth: 693 km\n"
    "Availability: Global, focusing on Europe.\n"
    "Year: Launched in 2014.\n"
    "Price: Free under the Copernicus Open License."
)

# Slide 5 - Sentinel-2
slide_5 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_5.shapes.title
content = slide_5.placeholders[1]

title.text = "#Sentinel-2 (Optical Satellite for Land Monitoring)"
content.text = (
    "Resolution: 10m, 20m, 60m\n"
    "Orbit Rotation: ~98.7 minutes (Sun-synchronous)\n"
    "Distance from Earth: 786 km\n"
    "Availability: Global.\n"
    "Year: Launched in 2015 (Sentinel-2A) and 2017 (Sentinel-2B).\n"
    "Price: Free."
)

# Slide 6 - Landsat 8/9
slide_6 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_6.shapes.title
content = slide_6.placeholders[1]

title.text = "#Landsat 8/9 (NASA Landsat Program)"
content.text = (
    "Resolution: 15m (panchromatic), 30m (multispectral)\n"
    "Orbit Rotation: ~99 minutes (Sun-synchronous)\n"
    "Distance from Earth: 705 km\n"
    "Availability: Global.\n"
    "Year: Landsat 8 (2013), Landsat 9 (2021).\n"
    "Price: Free."
)

# Slide 7 - ALOS World 3D (AW3D30)
slide_7 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_7.shapes.title
content = slide_7.placeholders[1]

title.text = "#ALOS World 3D (AW3D30)"
content.text = (
    "Resolution: 30m\n"
    "Orbit Rotation: ~100 minutes\n"
    "Distance from Earth: 691 km\n"
    "Availability: Global, via JAXA.\n"
    "Year: Launched in 2006.\n"
    "Price: Free."
)

# Slide 8 - TanDEM-X
slide_8 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_8.shapes.title
content = slide_8.placeholders[1]

title.text = "#TanDEM-X (German Aerospace Center)"
content.text = (
    "Resolution: 12m (High-res DEM)\n"
    "Orbit Rotation: ~97.6 minutes\n"
    "Distance from Earth: 514 km\n"
    "Availability: Global, commercial and scientific use.\n"
    "Year: Launched in 2010.\n"
    "Price: Commercial (varies by use case)."
)

# Slide 9 - ArcticDEM
slide_9 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_9.shapes.title
content = slide_9.placeholders[1]

title.text = "#ArcticDEM (Polar Geospatial Center)"
content.text = (
    "Resolution: 2m to 10m\n"
    "Distance from Earth: 450-700 km (multiple high-resolution commercial satellites)\n"
    "Availability: Arctic regions.\n"
    "Year: Continuously updated in 2024.\n"
    "Price: Free."
)

# Slide 10 - Commercial Satellite Data
slide_10 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_10.shapes.title
content = slide_10.placeholders[1]

title.text = "2024 Available Commercial Satellite Data"
content.text = (
    "#WorldView (Maxar Technologies)\n"
    "- Resolution: Up to 31cm (WorldView-3)\n"
    "- Price: ~$15-25 USD per km²\n\n"
    "#PlanetScope (Planet Labs)\n"
    "- Resolution: 3-5m\n"
    "- Price: ~$1,000 - $3,000 USD per 1,000 km²\n\n"
    "#SPOT 6/7 (Airbus Defence and Space)\n"
    "- Resolution: 1.5m\n"
    "- Price: €2-5 EUR per km²\n\n"
    "#Pleiades (Airbus Defence and Space)\n"
    "- Resolution: 50cm\n"
    "- Price: €20-30 EUR per km²\n\n"
    "#SkySat (Planet Labs)\n"
    "- Resolution: 50cm\n"
    "- Price: $50-100 USD per km²"
)

# Final Slide
slide_final = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_final.shapes.title
content = slide_final.placeholders[1]

title.text = "Conclusion"
content.text = (
    "In 2024, access to satellite imagery and DEM data continues to grow, offering a range of options for public and commercial use. "
    "From low Earth orbits to high-resolution satellites, there’s a wealth of data available to support geospatial projects worldwide."
)

# Save the presentation
ppt_file_path = "path/Geospatial_Data_Presentation.pptx"
ppt.save(ppt_file_path)

ppt_file_path

